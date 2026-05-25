import streamlit as st
from pypdf import PdfReader
from docx import Document
from pptx import Presentation
from PIL import Image
import pytesseract
import chromadb
from openai import OpenAI
from fpdf import FPDF
try:
    import speech_recognition as sr
    voice_enabled = True
except:
    voice_enabled = False
from sentence_transformers import SentenceTransformer
# ---------------- PAGE CONFIG ---------------- #

st.set_page_config(
    page_title="AI Multi Document Assistant",
    layout="wide"
)

# ---------------- CUSTOM CSS ---------------- #

st.markdown("""
<style>

body {
    background-color: #0f1117;
}

.main {
    background-color: #0f1117;
    color: white;
}

.stButton button {
    background-color: #16a34a;
    color: white;
    border-radius: 10px;
    padding: 10px;
    border: none;
    font-size: 16px;
}

.stChatInput input {
    border-radius: 12px;
}

</style>
""", unsafe_allow_html=True)

# ---------------- SIDEBAR ---------------- #

st.sidebar.title("💬 Chat History")

if "history" not in st.session_state:
    st.session_state.history = []

if st.sidebar.button("🆕 New Chat"):
    st.session_state.history = []

for chat in st.session_state.history:

    st.sidebar.markdown(f"""
<div style="
background-color:#1f2937;
padding:10px;
border-radius:10px;
margin-bottom:10px;
color:white;
font-size:14px;
">

📌 {chat['title']}

</div>
""", unsafe_allow_html=True)

# ---------------- TITLE ---------------- #

st.title("🧠 AI Knowledge Assistant using RAG")

st.write(
    "Upload PDF, DOCX, PPTX, Images and ask questions."
)

# ---------------- FILE UPLOAD ---------------- #

uploaded_files = st.file_uploader(
    "Upload Files",
    type=["pdf", "docx", "pptx", "png", "jpg", "jpeg"],
    accept_multiple_files=True
)

# ---------------- CHROMADB ---------------- #

client = chromadb.PersistentClient(
    path="./chroma_db"
)

try:
    client.delete_collection("documents")
except:
    pass

collection = client.create_collection(
    name="documents"
)
# ---------------- OPENROUTER API ---------------- #

client_ai = OpenAI(
    api_key=st.secrets["OPENROUTER_API_KEY"],
    base_url="https://openrouter.ai/api/v1"
)

# ---------------- TEXT EXTRACTION ---------------- #

def extract_pdf(file):

    text = ""

    pdf_reader = PdfReader(file)

    for page in pdf_reader.pages:

        extracted = page.extract_text()

        if extracted:
            text += extracted + "\n"

    return text


def extract_docx(file):

    doc = Document(file)

    text = ""

    for para in doc.paragraphs:
        text += para.text + "\n"

    return text


def extract_pptx(file):

    prs = Presentation(file)

    text = ""

    for slide in prs.slides:

        for shape in slide.shapes:

            if hasattr(shape, "text"):

                text += shape.text + "\n"

    return text


def extract_image(file):

    image = Image.open(file)

    text = pytesseract.image_to_string(image)

    return text

# ---------------- PROCESS FILES ---------------- #

all_text = ""

if uploaded_files:

    for uploaded_file in uploaded_files:

        st.write(
            f"📂 Processing: {uploaded_file.name}"
        )

        file_extension = uploaded_file.name.split(".")[-1].lower()

        text = ""

        if file_extension == "pdf":

            st.success("📄 PDF detected")

            text = extract_pdf(uploaded_file)

        elif file_extension == "docx":

            st.success("📝 DOCX detected")

            text = extract_docx(uploaded_file)

        elif file_extension == "pptx":

            st.success("📊 PPTX detected")

            text = extract_pptx(uploaded_file)

        elif file_extension in ["png", "jpg", "jpeg"]:

            st.success("🖼 Image detected")

            text = extract_image(uploaded_file)

        else:

            st.error("Unsupported file type")

            continue

        all_text += text + "\n"

    # ---------------- CHUNKING ---------------- #

    chunks = []

    chunk_size = 500

    for i in range(0, len(all_text), chunk_size):

        chunks.append(all_text[i:i + chunk_size])

    st.success(
        f"Total Chunks Created: {len(chunks)}"
    )

    # ---------------- EMBEDDINGS ---------------- #
    embeddings = embedding_model.encode(
        chunks).tolist()
    collection.add(
        ids=[str(i) for i in range(len(chunks))],
        embeddings=embeddings,
        documents=chunks
    )
    
    st.success("Files processed successfully ✅")

    # ---------------- FULL NOTES ---------------- #

    if st.button("📄 Generate Full Notes"):

        try:

            summary_prompt = f"""
You are an advanced AI teacher and educational assistant.

Read the uploaded document carefully and generate COMPLETE STUDY NOTES.

The notes should be VERY DETAILED, EASY TO UNDERSTAND, and useful for:

- School students
- College students
- Exam preparation
- Interview preparation

Generate the notes in this structure:

1. Introduction
2. Main Concepts
3. Definitions
4. Detailed Explanation
5. Key Points
6. Examples
7. Real-world Applications
8. Advantages and Disadvantages
9. Important Formulas
10. Revision Notes
11. Summary

Rules:
- Generate LONG and DETAILED notes
- Explain in simple language
- Use headings and bullet points
- Cover maximum concepts from the document

Document Content:
{all_text}
"""

            summary_response = client_ai.chat.completions.create(

                model="openai/gpt-3.5-turbo",

                messages=[
                    {
                        "role": "user",
                        "content": summary_prompt
                    }
                ]

            )

            summary = summary_response.choices[0].message.content

            st.markdown(f"""
<div style="
background-color:#111827;
padding:25px;
border-radius:15px;
color:white;
line-height:1.8;
font-size:16px;
margin-top:20px;
">

<h2 style="color:#22c55e;">📘 Full Study Notes</h2>

<div>

{summary}

</div>
</div>
""", unsafe_allow_html=True)

            # ---------------- PDF DOWNLOAD ---------------- #

            pdf = FPDF()

            pdf.add_page()

            pdf.set_font("Arial", size=12)

            pdf.multi_cell(0, 10, summary)

            pdf.output("study_notes.pdf")

            with open("study_notes.pdf", "rb") as file:

                st.download_button(
                    "⬇ Download Notes PDF",
                    file,
                    file_name="study_notes.pdf"
                )

        except Exception as e:

            st.error(f"Error: {e}")

# ---------------- VOICE INPUT ---------------- #

if voice_enabled:

    st.subheader("🎤 Voice Question")

    question = ""

    if st.button("🎙 Start Recording"):

        recognizer = sr.Recognizer()

        try:

            with sr.Microphone() as source:

                st.info("Speak now...")

                audio = recognizer.listen(source)

                voice_text = recognizer.recognize_google(audio)

                st.success(f"You said: {voice_text}")

                question = voice_text

        except Exception as e:

            st.error("Microphone not supported on deployment.")

else:

    st.warning("Voice feature disabled on cloud deployment.")

# ---------------- CHAT INPUT ---------------- #

user_question = st.chat_input(
    "Ask anything..."
)

if user_question:

    question = user_question

# ---------------- QUESTION ANSWERING ---------------- #

if question:

    st.chat_message("user").write(question)

    try:

        context = ""

        # SEARCH DOCUMENTS
        if uploaded_files and collection.count() > 0:

            results = collection.query(
                query_texts=[question],
                n_results=3
            )

            docs = results["documents"][0]

            context = "\n".join(docs)

        # DOCUMENT MODE
        if context.strip():

            prompt = f"""
You are an AI study assistant.

Use the document context to answer clearly.

Context:
{context}

Question:
{question}
"""

        # GENERAL AI MODE
        else:

            prompt = f"""
You are a helpful AI assistant.

Answer clearly and simply.

Question:
{question}
"""

        # ---------------- AI RESPONSE ---------------- #

        response = client_ai.chat.completions.create(

            model="openai/gpt-3.5-turbo",

            messages=[
                {
                    "role": "user",
                    "content": prompt
                }
            ]

        )

        answer = response.choices[
            0
        ].message.content.strip()

        st.chat_message(
            "assistant"
        ).write(answer)

        # ---------------- GENERATE CHAT TITLE ---------------- #

        title_prompt = f"""
Generate a short title for this chat.

Question:
{question}

Rules:
- Maximum 5 words
- Short and professional
- No quotes
"""

        title_response = client_ai.chat.completions.create(

            model="openai/gpt-3.5-turbo",

            messages=[
                {
                    "role": "user",
                    "content": title_prompt
                }
            ]

        )

        chat_title = title_response.choices[
            0
        ].message.content.strip()

        # ---------------- SAVE HISTORY ---------------- #

        chat_item = {
            "title": chat_title,
            "question": question,
            "answer": answer
        }

        st.session_state.history.append(
            chat_item
        )

    except Exception as e:

        st.error(f"Error: {e}")